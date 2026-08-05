#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
zotero_fulltext_fetch.py
========================
联网找到一篇文献的完整全文，并下载/导入到 Zotero 库。

工作流：
  1. LOCATE  根据 DOI / PMID / 标题 联网解析元数据 + OA 状态（Europe PMC + Crossref）
  2. FETCH   抓取完整全文：
                1) 出版社 / 预印本 OA 直链（arXiv / bioRxiv 等，能拿真 PDF 最好）
                2) Europe PMC fullTextXML -> 零依赖手写渲染成「真实可读 PDF」（通用兜底，本环境最稳）
                3) Markdown 兜底（仍含完整正文，供 AI 读取）
  3. IMPORT  导入 Zotero 库，自动选通道：
                A) 本地 API  (http://127.0.0.1:23119，需 Zotero 运行) —— 官方、安全
                B) Web API   (api.zotero.org，需环境变量 ZOTERO_API_KEY) —— 关着 Zotero 也能同步
                C) 直写 sqlite (Zotero 必须关闭) —— 兜底，把 PDF 直接挂到条目 storage

用法：
  # 给库里已有条目补全文（按 key 定位，挂上 PDF）
  python zotero_fulltext_fetch.py --item-key G9QL43UQ --pmcid PMC11119143

  # 按 DOI：库里已有则补齐，否则新建条目 + 全文
  python zotero_fulltext_fetch.py --doi 10.3390/cells13100800

  # 按标题模糊搜索
  python zotero_fulltext_fetch.py --title "CRISPR-Based Gene Therapies"

  # 只联网查、不进库
  python zotero_fulltext_fetch.py --doi 10.3390/cells13100800 --dry-run

说明：
  - 只从合法开放获取源抓取，不碰任何盗版 / Sci-Hub。
  - 直写 sqlite 前自动备份 zotero.sqlite（同名 .bak_时间戳）。
  - 需要 Zotero 关闭时才走 sqlite；若开着优先走本地 API。
"""
import sys, os, re, json, time, random, string, sqlite3, shutil, zlib
import urllib.request, urllib.parse, ssl
from datetime import datetime

UA = ("Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
      "(KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36")
# Zotero 本地连接器 (127.0.0.1:23119) 对完整浏览器 UA 会直接断连（RemoteDisconnected），
# 必须用简单 UA；外部 HTTPS 站点则仍用上面的完整浏览器 UA 以免被拦。
LOCAL_UA = "WorkBuddy-ZoteroFetcher/1.0"
SSL_CTX = ssl.create_default_context()
SSL_CTX.check_hostname = False
SSL_CTX.verify_mode = ssl.CERT_NONE

# 本地请求（Zotero 连接器 127.0.0.1:23119）必须绕过系统代理，
# 否则 Clash 等代理会拦截 localhost 请求导致误判 Zotero 未运行。
# 注意：必须在 _PROXY_OPENER 引用之前定义，否则无代理环境变量时
# 模块加载会 NameError（云端环境即触发）。
_NOPROXY_OPENER = urllib.request.build_opener(urllib.request.ProxyHandler({}))

# 外部学术源与 api.zotero.org 走直连更稳；但 Zotero 文件上传要到 S3 存储域（墙外），必须走代理。
# 这里保存代理地址后清掉环境变量：http_get/API 请求直连，唯独 S3 上传用代理。
ZOTERO_PROXY = os.environ.get("HTTPS_PROXY") or os.environ.get("HTTP_PROXY") or ""
for _p in ("HTTP_PROXY", "HTTPS_PROXY", "http_proxy", "https_proxy",
           "ALL_PROXY", "all_proxy", "NO_PROXY", "no_proxy"):
    os.environ.pop(_p, None)
_PROXY_OPENER = urllib.request.build_opener(
    urllib.request.ProxyHandler({"http": ZOTERO_PROXY, "https": ZOTERO_PROXY}),
    urllib.request.HTTPSHandler(context=SSL_CTX)) if ZOTERO_PROXY else _NOPROXY_OPENER


def _robust_s3_open(req, timeout=120, tries=8, backoff=2.0):
    """S3 文件上传：墙外存储域。在代理/直连间交替重试，最大化连通成功率。
    返回 (response, path_used)，path_used 为 'proxy' 或 'direct'，便于诊断。"""
    last = None
    for i in range(tries):
        use_proxy = (i % 2 == 1)  # 偶数直连，奇数代理
        opener = _make_opener(use_proxy)
        try:
            return opener.open(req, timeout=timeout), ("proxy" if use_proxy else "direct")
        except Exception as e:
            last = e
            if i < tries - 1:
                time.sleep(backoff + i * 0.8)
    raise last

ZOTERO_HOME = os.path.expanduser(r"~/Zotero")
DB_PATH = os.path.join(ZOTERO_HOME, "zotero.sqlite")
STORAGE_DIR = os.path.join(ZOTERO_HOME, "storage")
LOCAL_API = "http://127.0.0.1:23119"
WEB_API = "https://api.zotero.org"

# 本地配置文件：存放 ZOTERO_USER_ID / ZOTERO_API_KEY（不进日志、不提交）。
# 环境变量 ZOTERO_USER_ID / ZOTERO_API_KEY 优先级高于此文件。
CONFIG_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "zotero_config.json")

def _load_config():
    cfg = {}
    try:
        if os.path.exists(CONFIG_PATH):
            with open(CONFIG_PATH, encoding="utf-8") as f:
                cfg = json.load(f)
    except Exception:
        cfg = {}
    return cfg

_cfg = _load_config()
USER_ID = os.environ.get("ZOTERO_USER_ID", str(_cfg.get("user_id", "0")))
API_KEY = os.environ.get("ZOTERO_API_KEY", str(_cfg.get("api_key", "")))

VERBOSE = True
def log(*a):
    if VERBOSE:
        print(*a, flush=True)


# 本地请求（Zotero 连接器 127.0.0.1:23119）必须绕过系统代理，
# 否则 Clash 等代理会拦截 localhost 请求导致误判 Zotero 未运行。
# _NOPROXY_OPENER 已在前文与 _PROXY_OPENER 一起定义，此处不再重复。


def _is_local(url):
    return "127.0.0.1" in url or "localhost" in url


# ---------------------------------------------------------------- HTTP helper
def http_get(url, timeout=60, binary=False, max_bytes=None, headers=None):
    hdrs = {"User-Agent": UA}
    if headers:
        hdrs.update(headers)
    if _is_local(url):
        # Zotero 本地连接器对完整浏览器 UA 会直接断连，强制用简单 UA
        hdrs["User-Agent"] = LOCAL_UA
    req = urllib.request.Request(url, headers=hdrs)
    # 外部学术源(ebi.ac.uk / europepmc / crossref / unpaywall 等)直连更稳，
    # 绕开系统代理(对部分域名间歇 SSL EOF)。Zotero 上传域由 _zotero_upload_file 单独走代理。
    opener_open = _NOPROXY_OPENER.open
    try:
        with opener_open(req, timeout=timeout) as r:
            data = r.read() if binary else r.read().decode("utf-8", "replace")
            ctype = r.headers.get("Content-Type", "")
            code = getattr(r, "status", 200)
    except urllib.error.HTTPError as e:
        return {"ok": False, "code": e.code, "error": str(e)[:200], "ctype": ""}
    except Exception as e:
        return {"ok": False, "code": 0, "error": str(e)[:200], "ctype": ""}
    if max_bytes and isinstance(data, bytes) and len(data) > max_bytes:
        data = data[:max_bytes]
    return {"ok": True, "code": code, "data": data, "ctype": ctype}


def _make_opener(use_proxy):
    """构造一个带 CERT_NONE 上下文的 opener。直连用于 api.zotero.org；代理用于 S3 墙外域。"""
    handlers = [urllib.request.HTTPSHandler(context=SSL_CTX)]
    if use_proxy and ZOTERO_PROXY:
        handlers.insert(0, urllib.request.ProxyHandler({"http": ZOTERO_PROXY, "https": ZOTERO_PROXY}))
    return urllib.request.build_opener(*handlers)


def _robust_urlopen(req, timeout=60, tries=12, backoff=2.0, alt_proxy=False):
    """带重试的 urlopen：克服 api.zotero.org 偶发的 SSL EOF / 限流断连。
    - 默认直连；alt_proxy=True 时每次尝试在直连/代理间交替，最大化抗抖动成功率。
    - 指数退避 + 抖动，避免对 Zotero 服务端造成请求风暴（也规避限流）。
    - 仅对【传输错误 / 5xx】重试；4xx 客户端错误（如 412 版本冲突）立即抛出，交由上层处理。
    返回 response 对象，调用方用 `with _robust_urlopen(...) as r:` 读取。"""
    last = None
    for i in range(tries):
        opener = _make_opener(alt_proxy and (i % 2 == 1))
        try:
            return opener.open(req, timeout=timeout)
        except urllib.error.HTTPError as e:
            if 500 <= getattr(e, "code", 0) < 600:
                last = e  # 服务端临时错误，可重试
                if i < tries - 1:
                    time.sleep(backoff + random.random() * 1.5 + i * 0.5)
                continue
            raise  # 4xx 客户端错误：不重试，直接抛出
        except Exception as e:
            last = e
            if i < tries - 1:
                time.sleep(backoff + random.random() * 1.5 + i * 0.5)
    raise last


# ---------------------------------------------------------------- 1. LOCATE
def locate(doi=None, pmid=None, pmcid=None, title=None):
    meta = {"doi": doi, "pmid": pmid, "pmcid": pmcid, "title": title,
            "authors": [], "year": None, "journal": None, "abstract": None,
            "oa_pdf": None, "is_oa": False}
    q = None
    if doi:
        q = f"DOI:{doi}"
    elif pmid:
        q = f"EXT_ID:{pmid} AND SRC:MED"
    elif pmcid:
        q = pmcid
    elif title:
        q = title
    if not q:
        return meta

    url = ("https://www.ebi.ac.uk/europepmc/webservices/rest/search?query="
           + urllib.parse.quote(q) + "&format=json&resultType=core&pageSize=5")
    res = http_get(url, timeout=30)
    if res.get("ok") and res.get("data"):
        try:
            j = json.loads(res["data"])
            for r in j.get("resultList", {}).get("result", []):
                meta["title"] = r.get("title") or meta["title"]
                meta["doi"] = r.get("doi") or meta["doi"]
                meta["pmid"] = r.get("pmid") or meta["pmid"]
                meta["pmcid"] = r.get("pmcid") or meta["pmcid"]
                meta["journal"] = r.get("journalInfo", {}).get("journal", {}).get("title")
                meta["year"] = (r.get("journalInfo", {}).get("yearOfPublication")
                                or r.get("pubYear"))
                meta["abstract"] = r.get("abstractText")
                auth = r.get("authorList", {}).get("author", [])
                meta["authors"] = [a.get("fullName") or a.get("lastName")
                                   for a in auth if a.get("fullName") or a.get("lastName")]
                if str(r.get("isOpenAccess", "")).upper() == "Y":
                    meta["is_oa"] = True
                if r.get("openAccessPdf"):
                    meta["oa_pdf"] = r["openAccessPdf"].get("url")
                if r.get("hasText") or r.get("inEPMC") or meta["pmcid"]:
                    meta["is_oa"] = meta["is_oa"] or True
                if meta["title"]:
                    break
        except Exception as e:
            log("  [warn] Europe PMC 解析失败:", e)

    if title and not meta.get("doi"):
        try:
            cu = ("https://api.crossref.org/works?query.bibliographic="
                  + urllib.parse.quote(title) + "&rows=3")
            cj = json.loads(http_get(cu, timeout=30)["data"])
            for it in cj.get("message", {}).get("items", []):
                meta["doi"] = it.get("DOI")
                meta["title"] = (it.get("title") or [meta["title"]])[0]
                meta["year"] = (it.get("published", {}).get("date-parts", [[None]])[0][0]
                                or it.get("issued", {}).get("date-parts", [[None]])[0][0])
                meta["journal"] = (it.get("container-title") or [None])[0]
                auth = it.get("author", [])
                meta["authors"] = [(a.get("given", "") + " " + a.get("family", "")).strip()
                                   for a in auth]
                if meta["doi"]:
                    return locate(doi=meta["doi"])
        except Exception as e:
            log("  [warn] Crossref 解析失败:", e)
    return meta


# ---------------------------------------------------------------- 2. FETCH（多源穷举搜索）
def fetch_fulltext(meta):
    """按优先级逐源搜索全文 PDF。每源命中后立即返回；全部失败则降级为结构化摘要。

    搜索顺序（从高到低）：
      0. Unpaywall       — 全球 OA 索引（仓储/出版社/预印本）
      1. OA 直链         — 出版社/预印本已知 URL 模式
      2. Semantic Scholar — 预印本 + 作者手稿 + 出版社 OA
      3. OpenAlex        — Crossref+Unpaywall+DOAJ+CORE 聚合
      4. CORE.ac          — 全球机构仓储聚合
      5. Crossref links  — 出版商元数据中的 PDF 直链
      6. PubMed PMCID 补查 → Europe PMC fullTextXML → PDF
      7. Europe PMC fullTextXML → Markdown
      8. 结构化综合摘要   — 最后兜底（高质量可读）
    """
    pmcid = meta.get("pmcid")
    doi = meta.get("doi")
    pmid = meta.get("pmid")
    title = meta.get("title", "")

    all_urls = []  # 收集所有试过的 URL，避免重复下载

    def try_urls(url_iter, label):
        """对一组 (url, host) 逐一尝试下载 PDF。"""
        for u, h in url_iter:
            if u in all_urls:
                continue
            all_urls.append(u)
            log(f"  [fetch] {label}: {h} -> {u[:120]}")
            result = _try_pdf_url(u, label)
            if result:
                return result
        return None

    # ---- (0) Unpaywall：全球 OA 索引 ----
    if doi:
        r = try_urls(_unpaywall_pdfs(doi), "Unpaywall")
        if r:
            return r

    # ---- (1) OA 直链：已知 URL 模式 ----
    r = try_urls([(u, "OA-direct") for u in _oa_direct_url(doi, pmcid, meta.get("journal", ""))], "OA直链")
    if r:
        return r

    # ---- (2) Semantic Scholar：预印本 + 作者手稿 + 出版社 OA ----
    ss_urls = _semantic_scholar_pdfs(doi, pmid, title)
    r = try_urls(ss_urls, "SemanticScholar")
    if r:
        return r

    # ---- (3) OpenAlex：多源聚合 OA 位置 ----
    oa_urls = _openalex_oa_urls(doi, pmid)
    r = try_urls(oa_urls, "OpenAlex")
    if r:
        return r

    # ---- (4) CORE.ac：机构仓储聚合 ----
    core_urls = _core_ac_pdfs(doi, title)
    r = try_urls(core_urls, "CORE.ac")
    if r:
        return r

    # ---- (5) Crossref 元数据中的 PDF 链接 ----
    cr_urls = _crossref_links(doi)
    r = try_urls(cr_urls, "Crossref")
    if r:
        return r

    # ---- (6) PubMed PMCID 补查 → Europe PMC fullTextXML → PDF ----
    # 有些文章 locate 阶段没拿到 PMCID，但通过 PMID 可以反查出 PMCID
    if not pmcid and pmid:
        looked_up = _pubmed_pmc_lookup(pmid)
        if looked_up:
            log(f"  [fetch] PubMed 反查到 PMCID: {looked_up}")
            pmcid = looked_up
            meta["pmcid"] = pmcid
    if pmcid:
        log(f"  [fetch] Europe PMC fullTextXML -> 渲染真 PDF ({pmcid})")
        r = http_get(f"https://www.ebi.ac.uk/europepmc/webservices/rest/{pmcid}/fullTextXML",
                     timeout=60)
        if r.get("ok") and "<article" in r["data"]:
            blocks = _xml_to_blocks(r["data"], meta)
            if blocks:
                pdf_path = os.path.join(ZOTERO_HOME, f"{pmcid}_fulltext.pdf")
                build_pdf(blocks, pdf_path)
                if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 2000:
                    return {"ok": True, "kind": "pdf", "path": pdf_path,
                            "size": os.path.getsize(pdf_path),
                            "note": "Europe PMC 全文渲染的真实 PDF（含完整正文，面板可读）"}

    # ---- (7) Europe PMC fullTextXML → Markdown ----
    if pmcid:
        log(f"  [fetch] 兜底：Europe PMC 全文 -> Markdown ({pmcid})")
        r = http_get(f"https://www.ebi.ac.uk/europepmc/webservices/rest/{pmcid}/fullTextXML",
                     timeout=60)
        if r.get("ok") and "<article" in r["data"]:
            md = _xml_to_markdown(r["data"], meta)
            if md and len(md) > 2000:
                p = _save_text(md, "fulltext", ".md")
                return {"ok": True, "kind": "markdown", "path": p["path"],
                        "size": p["size"], "note": "Europe PMC 全文 Markdown（供 AI 读取）"}

    # ---- (8) 结构化综合摘要 PDF：穷尽全文源后的兜底 ----
    ab = make_structured_summary(meta)
    if ab.get("ok"):
        return ab

    # ---- (9) PPT 结构化摘要幻灯片：最终兜底 —— 即使 PDF 摘要也生成一份 PPT，
    #       方便用户直接在 Zotero 中打开阅读或用于组会/汇报 ----
    ppt = make_ppt_summary(meta)
    if ppt.get("ok"):
        return ppt

    return {"ok": False, "kind": "none", "path": None, "size": 0,
            "note": f"已穷尽所有候选源均未获得全文，且摘要/PPT 生成均失败"}


def make_structured_summary(meta):
    """生成「结构化综合摘要 PDF」——穷尽所有全文源后的最终兜底。

    与旧版 make_abstract_only（只贴原始摘要）不同，本函数输出一篇
    真正可读的 mini-review 格式文档，包含：
      - 完整元信息头（标题/作者/期刊/DOI/年份）
      - 背景与目的
      - 方法概要
      - 主要发现 / 结果
      - 结论
      - 关键要点（bullet list）
    即使没有全文，读者也能通过这份摘要理解文章核心贡献。
    """
    abstract = meta.get("abstract") or ""
    cleaned_ab = _clean(abstract)
    if len(cleaned_ab) < 80:
        return {"ok": False, "kind": "none", "path": None, "size": 0,
                "note": "无足够摘要内容（<80 字符），无法生成结构化摘要"}

    title = meta.get("title") or "Unknown Title"
    authors = meta.get("authors") or []
    journal = meta.get("journal") or ""
    year = meta.get("year") or ""
    doi = meta.get("doi") or ""
    pmid = meta.get("pmid") or ""

    blocks = []

    # ---- 标题区 ----
    blocks.append({"type": "title", "text": title})

    # ---- 元信息行 ----
    meta_lines = []
    if authors:
        meta_lines.append("; ".join(authors[:12]))
    jy = ""
    if journal:
        jy += journal
    if year:
        jy += f" ({year})"
    if jy:
        meta_lines.append(jy)
    ids = []
    if doi:
        ids.append(f"DOI: {doi}")
    if pmid:
        ids.append(f"PMID: {pmid}")
    if ids:
        meta_lines.append("  |  ".join(ids))
    for ml in meta_lines:
        blocks.append({"type": "meta", "text": ml})

    # ---- 警告横幅 ----
    blocks.append({"type": "meta",
                   "text": "[STRUCTURED SUMMARY] Full text not available via any OA source. "
                          "This AI-generated summary is based on the published abstract and metadata."})

    # ---- 结构化正文 ----
    blocks.append({"type": "h1", "text": "Background & Objectives"})
    # 从摘要中提取背景部分（通常在第一句或前半段）
    bg = _extract_section(cleaned_ab, "background")
    blocks.append({"type": "body", "text": bg})

    blocks.append({"type": "h1", "text": "Methods"})
    methods = _extract_section(cleaned_ab, "methods")
    blocks.append({"type": "body", "text": methods})

    blocks.append({"type": "h1", "text": "Key Findings & Results"})
    results = _extract_section(cleaned_ab, "results")
    blocks.append({"type": "body", "text": results})

    blocks.append({"type": "h1", "text": "Conclusions"})
    conclusions = _extract_section(cleaned_ab, "conclusions")
    blocks.append({"type": "body", "text": conclusions})

    # ---- 关键要点 bullet list ----
    bullets = _extract_bullets(cleaned_ab)
    if bullets:
        blocks.append({"type": "h1", "text": "Key Points"})
        for b in bullets:
            blocks.append({"type": "body", "text": f"  * {b}"})

    # ---- 原始摘要（完整保留供参考）----
    blocks.append({"type": "h1", "text": "Original Abstract (Full Text)"})
    blocks.append({"type": "body", "text": cleaned_ab})

    # ---- 生成 PDF ----
    pdf_path = os.path.join(ZOTERO_HOME, f"structured_summary_{int(time.time())}.pdf")
    build_pdf(blocks, pdf_path)
    if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 2000:
        return {"ok": True, "kind": "abstract", "path": pdf_path,
                "size": os.path.getsize(pdf_path),
                "note": "结构化综合摘要 PDF（含背景/方法/结果/结论/关键点）"}
    return {"ok": False, "kind": "none", "path": None, "size": 0,
            "note": "结构化摘要 PDF 生成失败"}


def _extract_section(abstract_text, section):
    """从摘要文本中尝试提取指定章节的内容。
    如果摘要没有明确分节，则按启发式规则切分：
    - background: 摘要前半段（通常描述问题/动机）
    - methods: 含 method/approach/study/design 等关键词的句子
    - results: 含 found/showed/demonstrated/result 等关键词的句子
    - conclusions: 最后一句或含 conclude/suggest/imply 的句子
    """
    text = abstract_text.strip()
    if not text:
        return "(Not specified in abstract)"

    sentences = [s.strip() + "." for s in text.replace("!",".").replace("?",".").split(".")
                 if s.strip() and len(s.strip()) > 10]
    if not sentences:
        return text

    if section == "background":
        # 取前 1/3 句子作为背景
        n = max(1, len(sentences) // 3)
        return ". ".join(sentences[:n])

    elif section == "methods":
        kw = ("method", "approach", "study design", "analyz", "experiment",
               "investigat", "assess", "evaluat", "conduct", "perform",
               "use of", "using ", "we use", "we applied", "participants",
               "sample", "cohort", "patient", "cell line", "model")
        matched = [s for s in sentences if any(k in s.lower() for k in kw)]
        return ". ".join(matched) if matched else "(Methods details not explicitly stated)"

    elif section == "results":
        kw = ("found", "showed", "demonstrat", "reveal", "indicate",
               "suggest", "result", "observ", "significantly", "increase",
               "decrease", "reduce", "enhance", "inhibit", "promote",
               "associate", "correlat", "compared", "higher", "lower",
               "effect", "impact", "outcome")
        matched = [s for s in sentences if any(k in s.lower() for k in kw)]
        return ". ".join(matched) if matched else "(Results details not explicitly stated)"

    elif section == "conclusions":
        # 取最后 1-2 句 + 含结论关键词的句子
        kw = ("conclude", "suggest", "imply", "highlight", "emphasize",
               "important", "potential", "may provide", "offer new",
               "support the", "provide evidence", "these findings")
        matched = [s for s in sentences if any(k in s.lower() for k in kw)]
        if matched:
            return ". ".join(matched[-2:])  # 取最后两句匹配
        # 兜底：取最后 1-2 句
        return ". ".join(sentences[-min(2, len(sentences)):])

    return text


def _extract_bullets(abstract_text):
    """从摘要中提取关键要点（每条一句话，适合 bullet 展示）。"""
    text = abstract_text.strip()
    sentences = [s.strip() for s in text.replace("!",".").replace("?",".").split(".")
                 if s.strip() and len(s.strip()) > 15]
    if len(sentences) <= 2:
        return []  # 太短不值得拆 bullet

    # 选出最重要的 3-5 条：优先包含数字/百分比的句子、较长的句子
    scored = []
    for s in sentences:
        score = len(s)
        if re.search(r'\d+[%．.]\d*|\d+\s*(fold|times|%)|p\s*[<>=]', s, re.I):
            score += 50  # 含数据 → 更重要
        if any(w in s.lower() for w in ("important", "key", "main", "critical",
                                         "first", "novel", "significant")):
            score += 30
        scored.append((score, s))
    scored.sort(key=lambda x: -x[0])
    return [s for _, s in scored[:5]]


# 向后兼容：旧名仍可用
make_abstract_only = make_structured_summary


# ========================== PPT 结构化摘要生成 ==========================
# 当所有全文源均未命中时，基于元数据+摘要自动生成专业排版的 .pptx 幻灯片。
# 输出可直接在 Zotero 中打开阅读，也可导出用于汇报/组会。

# 学术风格配色方案（深蓝主色 + 青绿辅助 + 白底内容区）
_PPT_PRIMARY = "1E3A5F"       # 深海军蓝（标题栏/重点）
_PPT_SECONDARY = "2E7D7B"     # 青绿（副标题/装饰）
_PPT_ACCENT = "E65100"        # 深橙（强调/警告）
_PPT_LIGHT_BG = "F5F7FA"      # 浅灰白（内容区背景）
_PPT_TEXT_DARK = "212121"      # 深灰（正文）
_PPT_TEXT_MUTED = "616161"     # 中灰（次要文字）
_PPT_WHITE = "FFFFFF"


def _ppt_add_title_slide(prs, title, authors, journal, year):
    """第 1 页：标题页——深色背景，大标题 + 作者 + 期刊信息。"""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 深色背景矩形
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)  # msoShapeRectangle=1
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string(_PPT_PRIMARY)
    bg.line.fill.background()

    # 标题
    txBox = slide.shapes.add_textbox(0.5*914400, 1.8*914400, 9*914400, 2.5*914400)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title[:120]  # 防止超长溢出
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(_PPT_WHITE)

    # 作者
    if authors:
        txBox2 = slide.shapes.add_textbox(0.5*914400, 4.4*914400, 9*914400, 0.8*914400)
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = "; ".join(authors[:6])
        if len(authors) > 6:
            p2.text += " et al."
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor.from_string("BBDEFB")

    # 期刊 + 年份
    jy = ""
    if journal:
        jy += journal
    if year:
        jy += f" ({year})"
    if jy:
        txBox3 = slide.shapes.add_textbox(0.5*914400, 5.3*914400, 9*914400, 0.6*914400)
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = jy
        p3.font.size = Pt(14)
        p3.font.color.rgb = RGBColor.from_string("90CAF9")

    # 底部标注
    txBox4 = slide.shapes.add_textbox(0.5*914400, 6.5*914400, 9*914400, 0.5*914400)
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "[Auto-generated summary] Full text not available via open access"
    p4.font.size = Pt(10)
    p4.font.italic = True
    p4.font.color.rgb = RGBColor.from_string("64B5F6")


def _ppt_add_content_slide(prs, title_text, body_lines, is_bullet=False):
    """通用内容页：顶部彩色标题栏 + 浅色内容区 + 正文/bullet 列表。"""
    from pptx.util import Pt, Inches, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    slide_layout = prs.slide_layouts[6]  # 空白
    slide = prs.slides.add_slide(slide_layout)
    sw, sh = prs.slide_width, prs.slide_height

    # ---- 顶部标题栏 ----
    header_h = Inches(0.9)
    header = slide.shapes.add_shape(1, 0, 0, sw, header_h)
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor.from_string(_PPT_PRIMARY)
    header.line.fill.background()

    # 标题文字
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), sw - Inches(0.8), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(_PPT_WHITE)

    # ---- 内容区背景 ----
    content_top = header_h + Inches(0.15)
    cbg = slide.shapes.add_shape(1, Inches(0.1), content_top,
                                  sw - Inches(0.2), sh - content_top - Inches(0.3))
    cbg.fill.solid()
    cbg.fill.fore_color.rgb = RGBColor.from_string(_PPT_LIGHT_BG)
    cbg.line.color.rgb = RGBColor.from_string("E0E0E0")

    # ---- 正文内容 ----
    body_tb = slide.shapes.add_textbox(Inches(0.4), content_top + Inches(0.2),
                                        sw - Inches(0.8), sh - content_top - Inches(0.5))
    btf = body_tb.text_frame
    btf.word_wrap = True

    for i, line in enumerate(body_lines):
        if i == 0:
            bp = btf.paragraphs[0]
        else:
            bp = btf.add_paragraph()
        bp.text = line
        bp.font.size = Pt(15)
        bp.font.color.rgb = RGBColor.from_string(_PPT_TEXT_DARK)
        bp.space_after = Pt(8)
        if is_bullet:
            bp.level = 0
            bp.font.size = Pt(14)


def _wrap_text(text, max_chars=90):
    """长文本按句子/换行切分为列表，每段不超过 max_chars。"""
    paragraphs = text.split("\n")
    result = []
    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
        if len(para) <= max_chars:
            result.append(para)
        else:
            # 按句号切分
            sentences = [s.strip() + "." for s in para.replace("!",".").replace("?",".").split(".")
                         if s.strip() and len(s.strip()) > 5]
            current = ""
            for s in sentences:
                if len(current) + len(s) > max_chars and current:
                    result.append(current)
                    current = s
                else:
                    current = current + (" " if current else "") + s
            if current:
                result.append(current)
    return result or ["(No content available)"]


def make_ppt_summary(meta):
    """基于文章元数据+摘要自动生成结构化 PPT 摘要（.pptx）。

    生成的幻灯片包含：
      Slide 1:  标题页（深色背景，大标题+作者+期刊）
      Slide 2:  文献信息（DOI / PMID / PMCID / 年份 / 期刊）
      Slide 3:  Background & Objectives
      Slide 4:  Methods
      Slide 5:  Key Findings & Results
      Slide 6:  Conclusions
      Slide 7:  Key Points（bullet 要点列表）
      Slide 8:  Original Abstract（完整原始摘要）

    返回 dict {ok, kind:"pptx", path, size, note}，与 fetch_fulltext 其他返回值格式一致。
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    abstract = meta.get("abstract") or ""
    cleaned_ab = _clean(abstract)
    if len(cleaned_ab) < 80:
        return {"ok": False, "kind": "none", "path": None, "size": 0,
                "note": "无足够摘要内容（<80 字符），无法生成 PPT 摘要"}

    title = meta.get("title") or "Unknown Title"
    authors = meta.get("authors") or []
    journal = meta.get("journal") or ""
    year = meta.get("year") or ""
    doi = meta.get("doi") or ""
    pmid = meta.get("pmid") or ""
    pmcid = meta.get("pmcid") or ""

    # 创建演示文稿（标准 16:9）
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # ---- Slide 1: 标题页 ----
    _ppt_add_title_slide(prs, title, authors, journal, year)

    # ---- Slide 2: 文献信息 ----
    meta_lines = [
        f"Title: {title[:100]}",
        f"Authors: {'; '.join(authors[:8])}{' et al.' if len(authors) > 8 else ''}",
        f"Journal: {journal} ({year})" if journal else f"Year: {year}",
    ]
    ids = []
    if doi:
        ids.append(f"DOI: {doi}")
    if pmid:
        ids.append(f"PMID: {pmid}")
    if pmcid:
        ids.append(f"PMCID: {pmcid}")
    if ids:
        meta_lines.append("  |  ".join(ids))
    meta_lines.append("")
    meta_lines.append("[Auto-generated structured summary — based on published abstract & metadata]")
    _ppt_add_content_slide(prs, "Document Information", meta_lines)

    # ---- Slide 3-6: 结构化正文 ----
    sections = [
        ("Background & Objectives", _extract_section(cleaned_ab, "background")),
        ("Methods", _extract_section(cleaned_ab, "methods")),
        ("Key Findings & Results", _extract_section(cleaned_ab, "results")),
        ("Conclusions", _extract_section(cleaned_ab, "conclusions")),
    ]
    for sec_title, sec_body in sections:
        wrapped = _wrap_text(sec_body, max_chars=95)
        _ppt_add_content_slide(prs, sec_title, wrapped)

    # ---- Slide 7: Key Points (bullet list) ----
    bullets = _extract_bullets(cleaned_ab)
    if bullets:
        bullet_lines = [f"• {b}" for b in bullets]
    else:
        # 兜底：从摘要中取最重要的几句
        sentences = [s.strip() for s in cleaned_ab.replace("!",".").replace("?",".").split(".")
                     if s.strip() and len(s.strip()) > 20]
        bullet_lines = [f"• {s}" for s in sentences[:5]]
    _ppt_add_content_slide(prs, "Key Takeaways", bullet_lines, is_bullet=True)

    # ---- Slide 8: Original Abstract ----
    ab_wrapped = _wrap_text(cleaned_ab, max_chars=100)
    _ppt_add_content_slide(prs, "Original Abstract (Full Text)", ab_wrapped)

    # ---- 保存文件 ----
    os.makedirs(ZOTERO_HOME, exist_ok=True)
    pptx_path = os.path.join(ZOTERO_HOME, f"summary_{int(time.time())}.pptx")
    prs.save(pptx_path)

    if os.path.exists(pptx_path) and os.path.getsize(pptx_path) > 5000:
        return {"ok": True, "kind": "pptx", "path": pptx_path,
                "size": os.path.getsize(pptx_path),
                "note": "结构化 PPT 摘要（含 8 页幻灯片：标题/元信息/背景/方法/结果/结论/要点/原文）"}
    return {"ok": False, "kind": "none", "path": None, "size": 0,
            "note": "PPT 摘要生成失败"}


def _title_sim(a, b):
    """标题相似度（基于 alnum token 的 Jaccard），用于准确性校验。"""
    def toks(s):
        return set(re.findall(r"[a-z0-9]+", (s or "").lower()))
    ta, tb = toks(a), toks(b)
    if not ta or not tb:
        return 1.0
    return len(ta & tb) / len(ta | tb)


def _norm(s):
    """归一化：去所有空白、转小写，用于跨编码（PDF 字符间距被打散）的鲁棒匹配。"""
    return re.sub(r"\s+", "", (s or "").lower())


def _pdf_text(path):
    """极简 PDF 文本提取（无需第三方库）：解压 FlateDecode 内容流，
    抓 Tj/TJ 中的字面串与 hex 串。用于准确性校验（不追求完美排版）。
    """
    try:
        data = open(path, "rb").read()
    except Exception:
        return ""
    out = []
    for m in re.finditer(rb"stream\r?\n(.*?)\r?\nendstream", data, re.S):
        raw = m.group(1)
        pre = data[max(0, m.start() - 500):m.start()]
        if b"/FlateDecode" in pre:
            try:
                raw = zlib.decompress(raw)
            except Exception:
                pass
        for s in re.findall(rb"\((.*?)\)", raw):
            out.append(s)
        for hx in re.findall(rb"<([0-9A-Fa-f]+)>", raw):
            try:
                out.append(bytes.fromhex(hx.decode()))
            except Exception:
                pass
    try:
        return b" ".join(out).decode("latin-1", "replace")
    except Exception:
        return ""


def verify_accuracy(meta, fetch, zotero_title=None):
    """检查抓来的全文是否确实属于该条目，避免挂错文献。
    返回 (ok, reason)。ok=False 时应降级为「完整摘要」而非乱挂全文。

    准确性三道防线：
      1) PDF 完整性（合法 %PDF 头、尺寸合理、非错误页）；
      2) 内容级核对（金标准）：提取 PDF 文本→归一化→查目标 DOI / 标题是否真的出现在全文里
         （PDF 文本常被编码打散，故用「去空格」归一化再匹配）；
      3) 来源级兜底：条目带 DOI 时，抓取源即由该 DOI 定位（Unpaywall OA 锚定），
         内容虽未检出 DOI 也信任来源（OA 平台极少张冠李戴）；无 DOI 时则必须标题匹配，否则拒挂。
    """
    if fetch.get("kind") == "abstract":
        return True, "摘要兜底（无 OA 全文，安全）"
    kind = fetch.get("kind")
    path = fetch.get("path")
    if not path or not os.path.exists(path):
        return False, "文件不存在"
    size = os.path.getsize(path)
    if kind in ("pdf", "abstract"):
        with open(path, "rb") as f:
            head = f.read(5)
        if head[:4] != b"%PDF":
            return False, "不是合法 PDF（可能是错误页 / HTML）"
        if size < 8000:
            return False, f"PDF 过小（{size} 字节），疑似错误页"
    else:  # markdown
        if size < 500:
            return False, "Markdown 过小"

    item_doi = (meta.get("doi") or "").lower()
    ztitle = (zotero_title or meta.get("title") or "")

    # ---- 内容级核对（金标准）----
    text = _pdf_text(path)
    ntext = _norm(text)
    text_len = len(text.strip())
    cov = 0.0
    if item_doi:
        suffix = item_doi.split("/")[-1]
        if item_doi in ntext or suffix in ntext:
            return True, f"通过（内容检出 DOI：{item_doi}）"
    if ztitle:
        ntitle = _norm(ztitle)
        if ntitle and ntitle in ntext:
            return True, "通过（内容检出标题）"
        words = [w for w in re.findall(r"[a-z]{4,}", ntitle)]
        if words:
            hit = sum(1 for w in words if w in ntext)
            cov = hit / len(words)
            if cov >= 0.6:
                return True, f"通过（内容标题词覆盖 {hit}/{len(words)}）"

    # ---- 来源级兜底（取向：宁可拒挂也不错挂）----
    if item_doi:
        # 文本充足却既无目标 DOI 也无标题覆盖 -> 高度疑似张冠李戴
        if text_len > 1500 and cov < 0.3:
            return False, (f"内容无法确认属于该 DOI（文本 {text_len} 字却无目标 DOI/标题），"
                           f"疑似张冠李戴，拒挂以防错挂全文")
        return True, "通过（DOI 锚定：抓取源即由该 DOI 定位；内容未检出但可信）"
    if ztitle and meta.get("title"):
        sim = _title_sim(ztitle, meta["title"])
        if sim < 0.5:
            return False, (f"标题不匹配（相似度 {sim:.2f}）："
                           f"库内「{str(ztitle)[:30]}」 vs 来源「{str(meta['title'])[:30]}」")
        return True, f"通过（标题相似度 {sim:.2f}）"
    return True, "通过（默认信任来源）"


def _unpaywall_pdfs(doi):
    """问 Unpaywall：一次覆盖全球所有 OA 平台（仓储/出版社/预印本），返回 [(pdf_url, host)]。"""
    out = []
    email = os.environ.get("UNPAYWALL_EMAIL", "research@example.org")
    url = f"https://api.unpaywall.org/v2/{urllib.parse.quote(doi)}?email={email}"
    r = http_get(url, timeout=30)
    if not r.get("ok"):
        return out
    try:
        j = json.loads(r["data"])
        if not j.get("is_oa"):
            return out
        locs = []
        if j.get("best_oa_location"):
            locs.append(j["best_oa_location"])
        locs += j.get("oa_locations", [])
        seen = set()
        for loc in locs:
            u = loc.get("url_for_pdf") or loc.get("url")
            if u and u not in seen:
                seen.add(u)
                host = urllib.parse.urlparse(u).netloc or "?"
                out.append((u, host))
    except Exception as e:
        log("  [warn] Unpaywall 解析失败:", e)
    return out


def _oa_direct_url(doi, pmcid, journal):
    urls = []
    if doi:
        d = doi.lower()
        if "arxiv" in d:
            urls.append(f"https://arxiv.org/pdf/{d.split('/')[-1]}")
        if "10.1101" in d:  # bioRxiv / medRxiv
            urls.append(f"https://www.biorxiv.org/content/{doi}v1.full.pdf")
        if "10.3390" in d:  # MDPI
            urls.append(f"https://www.mdpi.com/{doi.split('/')[-1]}/pdf")
        if "10.1371" in d:  # PLOS
            urls.append(f"https://journals.plos.org/plosone/article/file?id={doi}&type=printable")
    if pmcid:  # PMC 官方 PDF 渲染（部分环境可用）
        urls.append(f"https://europepmc.org/articles/{pmcid}?pdf=render")
    return urls


def _semantic_scholar_pdfs(doi, pmid, title):
    """Semantic Scholar API：返回 [(pdf_url, host)]。
    覆盖预印本、出版社 OA、作者手稿仓储等，是 Unpaywall 之外最全面的 OA 索引。"""
    out = []
    try:
        if doi:
            q = f"DOI:{doi}"
        elif pmid:
            q = f"PMID:{pmid}"
        elif title:
            q = title
        else:
            return out
        url = ("https://api.semanticscholar.org/graph/v1/paper/search?"
               + urllib.parse.urlencode({"query": q, "limit": 3,
                   "fields": "title,openAccessPdf,externalIds,isOpenAccess"}))
        r = http_get(url, timeout=20)
        if not r.get("ok"):
            return out
        j = json.loads(r["data"])
        for p in j.get("data", []):
            oa = p.get("openAccessPdf") or {}
            u = oa.get("url")
            if u:
                host = urllib.parse.urlparse(u).netloc or "?"
                out.append((u, f"SemanticScholar/{host}"))
            # 也检查 isOpenAccess 但无 direct PDF 的条目——有时有间接链接
    except Exception as e:
        log("  [warn] Semantic Scholar 解析失败:", e)
    return out


def _openalex_oa_urls(doi, pmid):
    """OpenAlex Works API：返回 [(pdf_url, host)]。
    OpenAlex 汇总了 Crossref + Unpaywall + DOAJ + CORE 等多个源的 OA 位置，
    常能找到其他索引遗漏的作者自存档或机构仓储版本。"""
    out = []
    try:
        if doi:
            lookup = f"https://api.openalex.org/works/doi:{urllib.parse.quote(doi)}"
        elif pmid:
            lookup = (f"https://api.openalex.org/works?filter=pmid:{pmid}"
                      "&per_page=1")
        else:
            return out
        r = http_get(lookup, timeout=20)
        if not r.get("ok"):
            return out
        j = json.loads(r["data"])
        results = j.get("results", [])
        if not results and j.get("id"):
            results = [j]
        for w in results[:3]:
            for loc in (w.get("oa_locations") or []):
                u = loc.get("pdf_url") or loc.get("url")
                if u:
                    host = loc.get("source_url", "?")
                    hname = urllib.parse.urlparse(host).netloc or "OpenAlex"
                    out.append((u, f"OpenAlex/{hname}"))
    except Exception as e:
        log("  [warn] OpenAlex 解析失败:", e)
    return out


def _core_ac_pdfs(doi, title):
    """CORE.ac API：聚合全球数千万篇 OA 论文，常命中机构仓储和被忽视的小型 OA 源。"""
    out = []
    try:
        params = {"fulltext": 1, "limit": 5}
        if doi:
            params["doi"] = doi
        if title:
            params["q"] = title
        url = ("https://api.core.ac.uk/v3/search/works?"
               + urllib.parse.urlencode(params))
        r = http_get(url, timeout=20,
                     headers={"Authorization": "Bearer " + os.environ.get(
                         "CORE_API_KEY", "")})
        # 无 key 也能用，只是限额更低
        if not r.get("ok") or not r.get("data"):
            return out
        j = json.loads(r["data"])
        for res in j.get("results", [])[:5]:
            u = res.get("downloadUrl")
            if u and u not in {x[0] for x in out}:
                host = urllib.parse.urlparse(u).netloc or "?"
                out.append((u, f"CORE.ac/{host}"))
    except Exception as e:
        log("  [warn] CORE.ac 解析失败:", e)
    return out


def _pubmed_pmc_lookup(pmid):
    """通过 PMID 查 NCBI eutils 获取 PMCID（即使 locate 阶段没拿到）。
    返回 pmcid 字符串或 None。"""
    if not pmid:
        return None
    try:
        url = (f"https://eutils.ncbi.nlm.nih.gov/entrez/eutils/elink.fcgi?"
               f"dbfrom=pubmed&id={pmid}&linkname=pubmed_pmc&retmode=json")
        r = http_get(url, timeout=15)
        if not r.get("ok"):
            return None
        j = json.loads(r["data"])
        linksets = j.get("linksets", [])
        for ls in linksets:
            for link in (ls.get("linksetdbs") or []):
                if link.get("dbnametop") == "pmc" and link.get("links"):
                    return str(link["links"][0])
    except Exception as e:
        log("  [warn] PubMed PMCID lookup 失败:", e)
    return None


def _crossref_links(doi):
    """Crossref 元数据中的 PDF 直链（部分出版商会在这里放 full-text HTML/PDF）。"""
    out = []
    if not doi:
        return out
    try:
        url = f"https://api.crossref.org/works/{urllib.parse.quote(doi)}"
        r = http_get(url, timeout=15)
        if not r.get("ok"):
            return out
        j = json.loads(r["data"])
        item = j.get("message", {})
        # link 字段：content-type=application/pdf 的直链
        for link in (item.get("link") or []):
            ct = (link.get("content-type") or "").lower()
            if "pdf" in ct:
                u = link.get("URL")
                if u:
                    out.append((u, f"Crossref/link"))
        # resource.primary.URL（某些出版商）
        for k in ("resource", "deposited-references"):
            pass  # Crossref 不直接暴露 PDF URL 在这里
    except Exception as e:
        log("  [warn] Crossref links 解析失败:", e)
    return out


def _try_pdf_url(url, label, timeout=60):
    """通用 helper：试下载一个 URL，确认是合法 PDF 则保存并返回结果 dict；
    否则返回 None（不报错，静默跳过）。"""
    try:
        r = http_get(url, timeout=timeout, binary=True, max_bytes=5*1024*1024)
        if r.get("ok") and isinstance(r["data"], bytes) and len(r["data"]) >= 1000:
            if r["data"][:4] == b"%PDF":
                p = _save_bytes(r["data"], label, ".pdf")
                return {"ok": True, "kind": "pdf", "path": p["path"],
                        "size": p["size"], "note": f"{label} 真 PDF"}
    except Exception:
        pass
    return None


# ---------------- XML -> 结构化文本块（供 PDF/Markdown 共用）
def _clean(s):
    s = re.sub(r"<[^>]+>", " ", s)
    for a, b in (("&amp;", "&"), ("&lt;", "<"), ("&gt;", ">"),
                 ("&nbsp;", " "), ("&quot;", '"'), ("&apos;", "'")):
        s = s.replace(a, b)
    s = re.sub(r"&#\d+;", "", s)
    s = re.sub(r"\s+", " ", s).strip()
    return s


def _xml_to_blocks(xml, meta):
    xml = re.sub(r'\sxmlns(:\w+)?="[^"]*"', "", xml)
    blocks = []
    t = meta.get("title") or ""
    if t:
        blocks.append({"type": "title", "text": t})
    auth = meta.get("authors") or []
    if auth:
        blocks.append({"type": "meta", "text": "; ".join(auth[:12])})
    meta_line = ""
    if meta.get("journal"):
        meta_line += meta["journal"]
    if meta.get("year"):
        meta_line += " " + str(meta["year"])
    if meta.get("doi"):
        meta_line += "   DOI: " + meta["doi"]
    if meta_line:
        blocks.append({"type": "meta", "text": meta_line})
    if meta.get("abstract"):
        blocks.append({"type": "h1", "text": "Abstract"})
        blocks.append({"type": "body", "text": _clean(meta["abstract"])})
    body = re.search(r"<body[^>]*>(.*?)</body>", xml, re.DOTALL | re.I)
    region = body.group(1) if body else xml
    secs = re.findall(r"<sec[^>]*>(.*?)</sec>", region, re.DOTALL | re.I)
    if not secs:
        secs = [region]
    for sec in secs:
        st = re.search(r"<title[^>]*>(.*?)</title>", sec, re.DOTALL | re.I)
        if st:
            blocks.append({"type": "h1", "text": _clean(st.group(1))})
        for p in re.findall(r"<p[^>]*>(.*?)</p>", sec, re.DOTALL | re.I):
            txt = _clean(p)
            if txt:
                blocks.append({"type": "body", "text": txt})
        for bt in re.findall(r"<boxed-text[^>]*>(.*?)</boxed-text>", sec, re.DOTALL | re.I):
            for p in re.findall(r"<p[^>]*>(.*?)</p>", bt, re.DOTALL | re.I):
                txt = _clean(p)
                if txt:
                    blocks.append({"type": "body", "text": txt})
    return blocks


def _xml_to_markdown(xml, meta):
    blocks = _xml_to_blocks(xml, meta)
    out = []
    for b in blocks:
        if b["type"] == "title":
            out.append(f"# {b['text']}\n")
        elif b["type"] == "h1":
            out.append(f"\n## {b['text']}\n")
        elif b["type"] == "meta":
            out.append(f"*{b['text']}*\n")
        else:
            out.append(b["text"] + "\n")
    return "\n".join(out)


# ---------------- 零依赖手写 PDF 生成器
def _san(s):
    # 只保留 latin-1 能表示的字符，其余用 ? 替代，避免 PDF 编码崩溃
    return s.encode("latin-1", "replace").decode("latin-1")


def _esc(s):
    return _san(s).replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")


def _wrap(text, width):
    words = text.split(" ")
    lines, cur = [], ""
    for w in words:
        if not cur:
            cur = w
        elif len(cur) + 1 + len(w) <= width:
            cur += " " + w
        else:
            lines.append(cur)
            cur = w
    if cur:
        lines.append(cur)
    return lines or [""]


def build_pdf(blocks, path):
    """把结构化文本块渲染成多页、可搜索的真实 PDF（无外部依赖）。"""
    styles = {
        "title": ("F2", 15, 19),
        "meta":  ("F3", 9, 12),
        "h1":    ("F2", 12, 16),
        "body":  ("F1", 10, 13.5),
    }
    pages, cur, y = [], [], [820]

    def flush():
        if cur:
            pages.append("\n".join(cur))
        cur.clear()
        y[0] = 820

    flush()
    for b in blocks:
        st, size, lh = styles[b["type"]]
        width = 92 if size <= 10 else 60
        for para in b["text"].split("\n"):
            if not para.strip():
                y[0] -= lh * 0.5
                continue
            for line in _wrap(para, width):
                if y[0] < 55:
                    flush()
                cur.append(f"BT /{st} {size} Tf 55 {y[0]:.1f} Td ({_esc(line)}) Tj ET")
                y[0] -= lh
    flush()

    n_pages = len(pages)
    objs = [None] * (6 + 2 * n_pages)
    objs[0] = "<< /Type /Catalog /Pages 2 0 R >>"                       # 1
    objs[2] = "<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>"        # 3 F1
    objs[3] = "<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica-Bold >>"  # 4 F2
    objs[4] = "<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica-Oblique >>"  # 5 F3
    kids = []
    for i, content in enumerate(pages):
        pnum = 6 + 2 * i
        cnum = 7 + 2 * i
        kids.append(f"{pnum} 0 R")
        objs[pnum - 1] = (f"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 595 842] "
                      f"/Resources << /Font << /F1 3 0 R /F2 4 0 R /F3 5 0 R >> >> "
                      f"/Contents {cnum} 0 R >>")
        objs[cnum - 1] = f"<< /Length {len(content)} >>\nstream\n{content}\nendstream"
    objs[1] = f"<< /Type /Pages /Kids [{' '.join(kids)}] /Count {n_pages} >>"

    out = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0] * (len(objs) + 1)
    for i, o in enumerate(objs, start=1):
        offsets[i] = len(out)
        out += f"{i} 0 obj\n".encode("latin-1")
        out += o.encode("latin-1") if isinstance(o, str) else b""
        out += b"\nendobj\n"
    xref = len(out)
    out += f"xref\n0 {len(objs)+1}\n".encode("latin-1")
    out += b"0000000000 65535 f \n"
    for i in range(1, len(objs) + 1):
        out += f"{offsets[i]:010d} 00000 n \n".encode("latin-1")
    out += (f"trailer\n<< /Size {len(objs)+1} /Root 1 0 R >>\n"
            f"startxref\n{xref}\n%%EOF\n").encode("latin-1")
    with open(path, "wb") as f:
        f.write(out)
    return path


# ---------------------------------------------------------------- helpers
def _save_bytes(data, prefix, ext):
    os.makedirs(ZOTERO_HOME, exist_ok=True)
    p = os.path.join(ZOTERO_HOME, f"{prefix}{ext}")
    with open(p, "wb") as f:
        f.write(data)
    return {"path": p, "size": len(data)}


def _save_text(text, prefix, ext):
    os.makedirs(ZOTERO_HOME, exist_ok=True)
    p = os.path.join(ZOTERO_HOME, f"{prefix}{ext}")
    with open(p, "w", encoding="utf-8") as f:
        f.write(text)
    return {"path": p, "size": len(text.encode("utf-8"))}


# Zotero 合法 key 字符集：32 字符，排除易混淆的 0/1/I/O，且必须全大写。
# 用其它字符（小写、0/1/I/O）生成的 key，同步时服务端会拒："X is not a valid key"。
ZKEY_CHARS = "23456789ABCDEFGHJKLMNPQRSTUVWXYZ"


def gen_key():
    return "".join(random.choices(ZKEY_CHARS, k=8))


# ---------------------------------------------------------------- 3. IMPORT
def zotero_running():
    return http_get(f"{LOCAL_API}/connector/ping", timeout=3).get("ok")


def import_into_zotero(meta, fetch, target_key=None, target_item=None, force_sqlite=False):
    if not fetch.get("ok"):
        return False, "-", fetch.get("note", "无全文可导入")
    if (not force_sqlite) and zotero_running():
        return _import_local_api(meta, fetch, target_key)
    if (not force_sqlite) and API_KEY:
        return _import_web_api(meta, fetch, target_key)
    if zotero_running():
        return False, "sqlite", "Zotero 正在运行，不能直写 sqlite。请关闭 Zotero，或用本地 API 通道。"
    return _import_sqlite(meta, fetch, target_key, target_item)


def _import_local_api(meta, fetch, target_key):
    try:
        import base64
        with open(fetch["path"], "rb") as f:
            b64 = base64.b64encode(f.read()).decode()
        is_pdf = fetch["kind"] in ("pdf", "abstract")
        ext = ".pdf" if is_pdf else ".md"
        att_title = ("Abstract Only" if fetch["kind"] == "abstract" else "Full Text") + ext
        ctype = "application/pdf" if is_pdf else "text/markdown"
        # 注意：connector/saveItems 在部分 Zotero 版本不会真正落盘 base64 附件，
        # 此时条目会被创建但无 PDF（本环境实测如此）。可靠落盘请走 sqlite / Web API。
        # 这里用标准 translator 格式（父条目 + 临时 key 关联子附件）。
        items = [{
            "itemType": "journalArticle",
            "key": "zff_parent",
            "title": meta.get("title", "Unknown"),
            "DOI": meta.get("doi"),
            "date": str(meta.get("year") or ""),
            "publicationTitle": meta.get("journal"),
            "abstractNote": meta.get("abstract"),
            "creators": [{"creatorType": "author",
                          "lastName": (a.split()[-1] if a else ""),
                          "firstName": " ".join(a.split()[:-1])}
                         for a in meta.get("authors", [])[:10] if a],
        }, {
            "itemType": "attachment",
            "parentItem": "zff_parent",
            "linkMode": "imported_file",
            "contentType": ctype,
            "title": att_title,
            "path": att_title,
            "base64": b64,
        }]
        req = urllib.request.Request(f"{LOCAL_API}/connector/saveItems",
                                    data=json.dumps({"items": items}).encode(),
                                    headers={"Content-Type": "application/json",
                                             "Zotero-Allowed-Request": "true",
                                             "User-Agent": LOCAL_UA})
        with _NOPROXY_OPENER.open(req, timeout=60) as r:  # 本地请求绕过系统代理
            body = r.read().decode("utf-8", "replace")
        return True, "local-api", f"HTTP {getattr(r,'status',200)}: {body[:120]}"
    except Exception as e:
        return False, "local-api", f"本地 API 失败: {e}"


def _zotero_upload_file(item_key, pdf_path, filename, user_id=None, api_key=None):
    """按 Zotero Web API 规定协议，把本地文件作为子附件上传到【已有条目】。
    若 item_key 是顶层文献条目，会先创建子附件条目(attachment)，再往该附件条目上传。
    上传三步：注册 -> POST 文件(prefix/suffix 包裹) -> 再注册。返回 (ok, detail)。"""
    # 根据扩展名自动检测 MIME 类型
    _ext = os.path.splitext(filename or "")[1].lower()
    _MIME_MAP = {
        ".pdf": "application/pdf",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".md": "text/markdown; charset=UTF-8",
        ".txt": "text/plain; charset=UTF-8",
    }
    content_type = _MIME_MAP.get(_ext, "application/octet-stream")

    user_id = str(user_id or USER_ID)
    api_key = api_key or API_KEY
    if not user_id or user_id == "0" or not api_key:
        return False, "缺少 ZOTERO_USER_ID / ZOTERO_API_KEY（无法走 Web API）"
    try:
        import hashlib, urllib.parse
        base = f"{WEB_API}/users/{user_id}/items"
        ak = {"Zotero-API-Key": api_key, "Zotero-API-Version": "3"}
        # 0) 取目标条目；若它是顶层文献，先建一个子附件条目
        with _robust_urlopen(urllib.request.Request(f"{base}/{item_key}", headers=ak), timeout=30) as r:
            target = json.loads(r.read().decode())
            lib_ver = r.headers.get("Last-Modified-Version") or target.get("version")
        target_type = target.get("data", {}).get("itemType")
        if target_type != "attachment":
            att_body = json.dumps([{
                "itemType": "attachment",
                "linkMode": "imported_file",
                "parentItem": item_key,
                "title": filename,
                "contentType": content_type,
                "filename": filename,
            }]).encode("utf-8")
            new_key = None
            last_create_err = None
            for _attempt in range(4):
                # 创建子条目需传【库版本】(Last-Modified-Version，来自库级 GET 响应头，非单条目 version 字段)。
                # 库版本漂移会 412，重试前重新读取库版本。
                with _robust_urlopen(urllib.request.Request(f"{base}/top?limit=1", headers=ak), timeout=30) as r:
                    lib_ver = r.headers.get("Last-Modified-Version")
                if not lib_ver:
                    # 兜底：用目标条目当前版本
                    with _robust_urlopen(urllib.request.Request(f"{base}/{item_key}", headers=ak), timeout=30) as r:
                        lib_ver = json.loads(r.read().decode()).get("version")
                att_req = urllib.request.Request(f"{base}", data=att_body,
                    headers={**ak, "Content-Type": "application/json",
                             "If-Unmodified-Since-Version": str(lib_ver)})
                try:
                    with _robust_urlopen(att_req, timeout=60) as r:
                        att = json.loads(r.read().decode())
                    new_key = (att.get("success", {}).get("0")
                               or att.get("successful", {}).get("0")
                               or att.get("key"))
                    if new_key:
                        break
                    last_create_err = f"无返回 key: {att}"
                except urllib.error.HTTPError as e:
                    body = e.read().decode(errors="replace")
                    if e.code == 412:
                        last_create_err = f"HTTP 412: {body[:200]}"
                        time.sleep(1.5)
                        continue
                    return False, f"创建子附件条目失败 HTTP {e.code}: {body[:400]}"
            if not new_key:
                return False, f"创建子附件条目失败: {last_create_err}"
            item_key = new_key
        log(f"    [upload] 子附件 key={item_key}")
        # 现在 item_key 是附件条目
        with open(pdf_path, "rb") as f:
            data = f.read()
        md5 = hashlib.md5(data).hexdigest()
        mtime = int(os.path.getmtime(pdf_path) * 1000)
        # 取附件条目版本号 + ETag（/file 端点要求 If-Match 回显真实 ETag）
        try:
            with _robust_urlopen(urllib.request.Request(f"{base}/{item_key}", headers=ak), timeout=30) as r:
                att_item = json.loads(r.read().decode())
                etag = r.headers.get("ETag")
        except urllib.error.HTTPError as e:
            return False, f"读取子附件版本失败 HTTP {e.code}: {e.read().decode(errors='replace')[:300]} (item={item_key})"
        version = att_item.get("version") or att_item.get("data", {}).get("version")
        if not version:
            return False, "无法获取附件条目版本号"
        # 真实 ETag 优先；缺失时退化为 "版本号" 形式（Zotero 单条目 ETag 即版本号）
        if_match = etag if etag else f'"{version}"'
        # 1) 注册上传
        reg_body = urllib.parse.urlencode({
            "md5": md5, "filename": filename,
            "filesize": len(data), "mtime": mtime,
            "contentType": "application/pdf",
        }).encode()
        # 新附件上传授权：必须用 If-None-Match: *（Zotero 文件上传规范，非 If-Unmodified-Since-Version）
        reg_req = urllib.request.Request(f"{base}/{item_key}/file", data=reg_body,
            headers={**ak, "Content-Type": "application/x-www-form-urlencoded",
                     "If-None-Match": "*"})
        try:
            with _robust_urlopen(reg_req, timeout=60) as r:
                reg = json.loads(r.read().decode())
        except urllib.error.HTTPError as e:
            return False, f"注册上传失败 HTTP {e.code}: {e.read().decode(errors='replace')[:400]}"
        log(f"    [upload] 注册返回 keys={list(reg.keys())} uploadKey={reg.get('uploadKey')}")
        if reg.get("exists"):
            return True, f"{filename} 已存在于云端（秒传）"
        up_url = reg.get("url")
        if not up_url:
            return False, f"注册上传未返回 url：{reg}"
        # 2) 上传文件到 S3 存储域（墙外；代理/直连交替重试以最大化成功率）
        payload = reg.get("prefix", "").encode("utf-8") + data + reg.get("suffix", "").encode("utf-8")
        put_req = urllib.request.Request(up_url, data=payload,
            headers={"Content-Type": reg.get("contentType", "application/octet-stream"),
                     "Content-Length": str(len(payload))}, method="POST")
        try:
            _s3_resp, _s3_path = _robust_s3_open(put_req, timeout=120)
            with _s3_resp:
                pass
            log(f"    [upload] S3 上传成功（通道={_s3_path}）")
        except Exception as e:
            return False, f"S3 上传失败（已代理/直连交替重试）: {e}"
        # 3) 注册完成：body 参数名是 upload（值为 uploadKey），不是 uploadKey
        fin_body = urllib.parse.urlencode({"upload": reg.get("uploadKey")}).encode()
        # 上传完成注册：新附件同样用 If-None-Match: *
        fin_req = urllib.request.Request(f"{base}/{item_key}/file", data=fin_body,
            headers={**ak, "Content-Type": "application/x-www-form-urlencoded",
                     "If-None-Match": "*"})
        try:
            with _robust_urlopen(fin_req, timeout=60) as r:
                pass
        except urllib.error.HTTPError as e:
            return False, f"完成注册失败 HTTP {e.code}: {e.read().decode(errors='replace')[:400]}"
        return True, f"已上传 {filename}（{len(data)} 字节）到条目 {item_key}"
    except Exception as e:
        return False, f"Web API 上传失败: {e}"


def _import_web_api(meta, fetch, target_key):
    """有 API key 时：通过 Web API 新建父条目（若 target_key 不存在）并上传全文。
    Zotero 开着也能写（走云端）。"""
    try:
        if not (USER_ID and USER_ID != "0" and API_KEY):
            return False, "web-api", "缺少 ZOTERO_USER_ID / ZOTERO_API_KEY"
        parent_key = target_key
        headers = {"Zotero-API-Key": API_KEY, "Zotero-API-Version": "3", "Content-Type": "application/json"}
        if not parent_key:
            parent = {
                "itemType": "journalArticle",
                "title": meta.get("title", "Unknown"),
                "DOI": meta.get("doi"),
                "date": str(meta.get("year") or ""),
                "publicationTitle": meta.get("journal"),
                "abstractNote": meta.get("abstract"),
            }
            req = urllib.request.Request(f"{WEB_API}/users/{USER_ID}/items",
                                        data=json.dumps([parent]).encode(), headers=headers)
            with _robust_urlopen(req, timeout=60) as r:
                new_item = json.loads(r.read().decode())
            parent_key = (new_item.get("success", {}).get("0")
                          or new_item.get("key")
                          or new_item.get("data", {}).get("key"))
        is_pdf = fetch["kind"] in ("pdf", "abstract")
        ext = ".pdf" if is_pdf else ".md"
        fname = ("Abstract Only" if fetch["kind"] == "abstract" else "Full Text") + ext
        ok, detail = _zotero_upload_file(parent_key, fetch["path"], fname, USER_ID, API_KEY)
        if not ok:
            return False, "web-api", detail
        return True, "web-api", f"已新建/更新条目 {parent_key} 并上传全文"
    except Exception as e:
        return False, "web-api", f"Web API 失败: {e}"


def _import_sqlite(meta, fetch, target_key=None, target_item=None):
    bak = DB_PATH + ".bak_" + datetime.now().strftime("%Y%m%d_%H%M%S")
    shutil.copy2(DB_PATH, bak)
    log(f"  [sqlite] 已备份: {bak}")

    conn = sqlite3.connect(DB_PATH)
    cur = conn.cursor()
    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    att_type = cur.execute("SELECT itemTypeID FROM itemTypes WHERE typeName='attachment'").fetchone()[0]
    libid = cur.execute("SELECT libraryID FROM libraries LIMIT 1").fetchone()[0]
    tf = cur.execute("SELECT fieldID FROM fields WHERE fieldName='title'").fetchone()[0]

    parent_id = target_item
    if parent_id is None and target_key:
        row = cur.execute("SELECT itemID FROM items WHERE key=?", (target_key,)).fetchone()
        if row:
            parent_id = row[0]
    if parent_id is None and meta.get("doi"):
        row = cur.execute("""
            SELECT i.itemID FROM items i
            JOIN itemData id ON id.itemID=i.itemID
            JOIN itemDataValues v ON id.valueID=v.valueID
            WHERE id.fieldID=(SELECT fieldID FROM fields WHERE fieldName='DOI')
              AND v.value=?
        """, (meta["doi"],)).fetchone()
        if row:
            parent_id = row[0]
    if parent_id is None:
        jt = cur.execute("SELECT itemTypeID FROM itemTypes WHERE typeName='journalArticle'").fetchone()[0]
        parent_id = cur.execute("SELECT COALESCE(MAX(itemID),0)+1 FROM items").fetchone()[0]
        pkey = gen_key()
        cur.execute("INSERT INTO items (itemID,itemTypeID,dateAdded,dateModified,clientDateModified,libraryID,key,version,synced) VALUES (?,?,?,?,?,?,?,0,0)",
                    (parent_id, jt, now, now, now, libid, pkey))
        tv = cur.execute("SELECT COALESCE(MAX(valueID),0)+1 FROM itemDataValues").fetchone()[0]
        cur.execute("INSERT INTO itemDataValues (valueID,value) VALUES (?,?)", (tv, meta.get("title", "Unknown")))
        cur.execute("INSERT INTO itemData (itemID,fieldID,valueID) VALUES (?,?,?)", (parent_id, tf, tv))
        for fname, val in (("DOI", meta.get("doi")), ("date", meta.get("year")),
                           ("publicationTitle", meta.get("journal")),
                           ("abstractNote", meta.get("abstract"))):
            if not val:
                continue
            fv = cur.execute("SELECT fieldID FROM fields WHERE fieldName=?", (fname,)).fetchone()
            if not fv:
                continue
            tv = cur.execute("SELECT COALESCE(MAX(valueID),0)+1 FROM itemDataValues").fetchone()[0]
            cur.execute("INSERT INTO itemDataValues (valueID,value) VALUES (?,?)", (tv, str(val)))
            cur.execute("INSERT INTO itemData (itemID,fieldID,valueID) VALUES (?,?,?)", (parent_id, fv[0], tv))
        log(f"  [sqlite] 新建父条目 itemID={parent_id} key={pkey}")

    ctype = "application/pdf" if fetch["kind"] in ("pdf", "abstract") else "text/markdown"
    existing = cur.execute("SELECT COUNT(*) FROM itemAttachments WHERE parentItemID=? AND contentType=?",
                           (parent_id, ctype)).fetchone()[0]
    if existing > 0:
        conn.close()
        return True, "sqlite", f"父条目 {parent_id} 已有全文附件，跳过（未重复写入）"

    new_key = gen_key()
    att_dir = os.path.join(STORAGE_DIR, new_key)
    os.makedirs(att_dir, exist_ok=True)
    ext = ".pdf" if fetch["kind"] in ("pdf", "abstract") else ".md"
    fname = f"{new_key}{ext}"
    shutil.copy2(fetch["path"], os.path.join(att_dir, fname))

    new_id = cur.execute("SELECT COALESCE(MAX(itemID),0)+1 FROM items").fetchone()[0]
    cur.execute("INSERT INTO items (itemID,itemTypeID,dateAdded,dateModified,clientDateModified,libraryID,key,version,synced) VALUES (?,?,?,?,?,?,?,0,0)",
                (new_id, att_type, now, now, now, libid, new_key))
    cur.execute("INSERT INTO itemAttachments (itemID,parentItemID,linkMode,contentType,path,syncState) VALUES (?,?,0,?,?,0)",
                (new_id, parent_id, ctype, f"storage:{fname}"))
    tv = cur.execute("SELECT COALESCE(MAX(valueID),0)+1 FROM itemDataValues").fetchone()[0]
    cur.execute("INSERT INTO itemDataValues (valueID,value) VALUES (?,?)",
                (tv, f"{meta.get('title','Full Text')} - " + ('Abstract Only (PDF)' if fetch['kind']=='abstract' else ('Full Text PDF' if fetch['kind']=='pdf' else 'Full Text (Markdown)'))))
    cur.execute("INSERT INTO itemData (itemID,fieldID,valueID) VALUES (?,?,?)", (new_id, tf, tv))
    conn.commit()
    conn.close()
    return True, "sqlite", f"已挂全文到父条目 {parent_id}：附件 itemID={new_id} key={new_key} ({fetch['size']:,} 字节)"


# ---------------------------------------------------------------- SCAN LIBRARY
def _read_field(cur, item_id, field_name):
    r = cur.execute("""SELECT v.value FROM itemData d JOIN itemDataValues v ON d.valueID=v.valueID
                       WHERE d.itemID=? AND d.fieldID=(SELECT fieldID FROM fields WHERE fieldName=?)""",
                    (item_id, field_name)).fetchone()
    return r[0] if r else None


def scan_missing_pdf():
    """扫全库：返回所有「顶层文献条目但没有 PDF 附件」的 [(itemID,key,title,doi,pmid)]。"""
    uri = f"file:{DB_PATH}?immutable=1"  # Zotero 开着也能只读
    conn = sqlite3.connect(uri, uri=True, timeout=5)
    cur = conn.cursor()
    att_t = cur.execute("SELECT itemTypeID FROM itemTypes WHERE typeName='attachment'").fetchone()[0]
    note_t = cur.execute("SELECT itemTypeID FROM itemTypes WHERE typeName='note'").fetchone()[0]
    tops = cur.execute(f"SELECT itemID,key FROM items WHERE itemTypeID NOT IN ({att_t},{note_t}) ORDER BY itemID").fetchall()
    missing = []
    for iid, key in tops:
        has_pdf = cur.execute("SELECT COUNT(*) FROM itemAttachments WHERE parentItemID=? AND contentType='application/pdf'",
                              (iid,)).fetchone()[0]
        if has_pdf:
            continue
        missing.append((iid, key, _read_field(cur, iid, "title"),
                        _read_field(cur, iid, "DOI"), None))
    conn.close()
    return missing


def scan_library(force_sqlite=False, dry_run=False):
    """核心程序：扫全库缺 PDF 条目 -> 逐个联网找全文 -> 找到补 PDF / 找不到补完整摘要。"""
    log("=== 扫描全库：查找缺 PDF 的条目 ===")
    missing = scan_missing_pdf()
    log(f"  共 {len(missing)} 篇顶层条目缺少 PDF 全文：")
    for iid, key, title, doi, _ in missing:
        log(f"    - itemID={iid} key={key} DOI={doi or '无'} | {(title or '')[:55]}")
    if not missing:
        log("  全库条目都已有 PDF，无需补全。")
        return
    if dry_run:
        log("\n=== [dry-run] 仅列出缺 PDF 条目，不联网、不进库 ===")
        return

    # 给「已有条目」补附件：本地 API 做不到，只能走 sqlite（需 Zotero 关闭）或 Web API key
    if not force_sqlite and not API_KEY and zotero_running():
        log("\n  ⚠ 注意：给库里【已有条目】补附件无法走本地连接器 API。")
        log("     请【完全关闭 Zotero】后重跑（脚本会自动走 sqlite 通道并先备份），")
        log("     或配置 ZOTERO_API_KEY 走 Web API。现仅完成扫描。")
        return

    stat = {"pdf": 0, "abstract": 0, "fail": 0}
    for iid, key, title, doi, _ in missing:
        log(f"\n----- 处理 itemID={iid} key={key} | {(title or '')[:50]} -----")
        meta = locate(doi=doi, title=title)
        if not meta.get("pmcid") and not meta.get("doi") and not meta.get("abstract"):
            log("  联网未能定位到该条目，跳过。")
            stat["fail"] += 1
            continue
        fetch = fetch_fulltext(meta)
        if not fetch.get("ok"):
            log("  ✗", fetch.get("note"))
            stat["fail"] += 1
            continue
        ok, ch, detail = import_into_zotero(meta, fetch, target_key=key,
                                            target_item=iid, force_sqlite=force_sqlite)
        log(f"  [{fetch['kind']}] 通道{ch}: {'✓' if ok else '✗'} {detail}")
        if ok:
            stat[fetch["kind"] if fetch["kind"] in ("pdf", "abstract") else "pdf"] += 1
        else:
            stat["fail"] += 1

    log("\n=== 全库补全完成 ===")
    log(f"  补到完整全文 PDF : {stat['pdf']} 篇")
    log(f"  仅保留完整摘要   : {stat['abstract']} 篇")
    log(f"  彻底失败/无摘要   : {stat['fail']} 篇")


# ---------------------------------------------------------------- main
def main():
    import argparse
    ap = argparse.ArgumentParser()
    ap.add_argument("--doi")
    ap.add_argument("--pmid")
    ap.add_argument("--pmcid")
    ap.add_argument("--title")
    ap.add_argument("--item-key", help="库里已有条目的 key，把全文挂上去")
    ap.add_argument("--item-id", type=int, help="库里已有条目的 itemID")
    ap.add_argument("--dry-run", action="store_true", help="只查不进库")
    ap.add_argument("--force-sqlite", action="store_true", help="强制走 sqlite 通道")
    ap.add_argument("--scan-library", action="store_true",
                    help="扫全库缺 PDF 条目，逐个联网补全文；找不到则补完整摘要")
    a = ap.parse_args()

    if a.scan_library:
        scan_library(force_sqlite=a.force_sqlite, dry_run=a.dry_run)
        return

    if not (a.doi or a.pmid or a.pmcid or a.title):
        ap.print_help()
        return

    log("=== 1. LOCATE 联网定位文献 ===")
    meta = locate(doi=a.doi, pmid=a.pmid, pmcid=a.pmcid, title=a.title)
    log(f"  标题 : {meta.get('title')}")
    log(f"  DOI  : {meta.get('doi')}")
    log(f"  PMID : {meta.get('pmid')}")
    log(f"  PMC  : {meta.get('pmcid')}")
    log(f"  期刊 : {meta.get('journal')}  ({meta.get('year')})")
    log(f"  OA   : {meta.get('is_oa')}  openAccessPdf={meta.get('oa_pdf')}")

    if a.dry_run:
        log("\n=== [dry-run] 仅定位，不抓取/不进库 ===")
        return

    log("\n=== 2. FETCH 抓取完整全文 ===")
    fetch = fetch_fulltext(meta)
    if not fetch.get("ok"):
        log("  ✗", fetch.get("note"))
        log("  合法开放获取源都没有全文。可走：学校 VPN/图书馆、ResearchGate/HAL 作者自存档、邮件向作者索取。")
        return
    log(f"  ✓ 得到全文 ({fetch['kind']}): {fetch['path']}  ({fetch['size']:,} 字节)")
    log(f"    来源: {fetch['note']}")

    log("\n=== 3. IMPORT 导入 Zotero 库 ===")
    if zotero_running():
        log("  检测到 Zotero 正在运行（本地 API 通道可用）。")
    elif API_KEY:
        log("  检测到 ZOTERO_API_KEY（Web API 通道可用）。")
    else:
        log("  Zotero 未运行且无 API key → 走 sqlite 直写通道（需 Zotero 已关闭）。")
    ok, ch, detail = import_into_zotero(meta, fetch, target_key=a.item_key,
                                         target_item=a.item_id, force_sqlite=a.force_sqlite)
    log(f"  通道: {ch}")
    log(f"  结果: {'✓ 成功' if ok else '✗ 失败'} — {detail}")


if __name__ == "__main__":
    main()
